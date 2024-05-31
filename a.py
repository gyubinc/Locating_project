import pptx
from pptx import Presentation
from pptx.util import Inches

# Create a presentation object
prs = Presentation()

# Slide 1: Title Slide
slide_1 = prs.slides.add_slide(prs.slide_layouts[0])
title = slide_1.shapes.title
subtitle = slide_1.placeholders[1]

title.text = "Dissecting Recall of Factual Association in Auto-Regressive Language Models"
subtitle.text = "Authors: Mor Geva et al.\nCategory: Knowledge Awareness, Knowledge Editing\nDate: April 28, 2023"

# Slide 2: Introduction
slide_2 = prs.slides.add_slide(prs.slide_layouts[1])
title = slide_2.shapes.title
content = slide_2.placeholders[1]

title.text = "Introduction"
content.text = "Transformer-based LMs store factual knowledge in parameters.\n\nResearch Goal: Study how the model gathers information to predict attributes for subject-relation queries using lens of information flow."

# Slide 3: Background
slide_3 = prs.slides.add_slide(prs.slide_layouts[1])
title = slide_3.shapes.title
content = slide_3.placeholders[1]

title.text = "Background"
content.text = "Focus on AR decoder-only models.\n\nStructure and function of Transformer-based models."

# Slide 4: Methodology
slide_4 = prs.slides.add_slide(prs.slide_layouts[1])
title = slide_4.shapes.title
content = slide_4.placeholders[1]

title.text = "Methodology"
content.text = "Basic information extraction setting.\n\nCritical points analysis: relation positions, subject positions.\n\n3-step mechanism of attribute extraction."

# Slide 5: Experimental Setup
slide_5 = prs.slides.add_slide(prs.slide_layouts[1])
title = slide_5.shapes.title
content = slide_5.placeholders[1]

title.text = "Experimental Setup"
content.text = "Data: COUNTERFACT dataset.\n\nModels: GPT-2, GPT-J."

# Slide 6: Results
slide_6 = prs.slides.add_slide(prs.slide_layouts[1])
title = slide_6.shapes.title
content = slide_6.placeholders[1]

title.text = "Results"
content.text = "Used attention blocking method.\n\nAnalyzed changes in relation and subject representations."

# Slide 7: Intermediate Subject Representations
slide_7 = prs.slides.add_slide(prs.slide_layouts[1])
title = slide_7.shapes.title
content = slide_7.placeholders[1]

title.text = "Intermediate Subject Representations"
content.text = "Analyzed hidden representations.\n\nSubject enrichment process and attribute extraction operation."

# Slide 8: Localizing Information Flow
slide_8 = prs.slides.add_slide(prs.slide_layouts[1])
title = slide_8.shapes.title
content = slide_8.placeholders[1]

title.text = "Localizing Information Flow"
content.text = "Described Attention Knockout methodology.\n\nAnalyzed experimental results."

# Slide 9: Subject Representation Analysis
slide_9 = prs.slides.add_slide(prs.slide_layouts[1])
title = slide_9.shapes.title
content = slide_9.placeholders[1]

title.text = "Subject Representation Analysis"
content.text = "Measured attribute rate.\n\nAnalyzed static embeddings and sublayers."

# Slide 10: Attribute Extraction
slide_10 = prs.slides.add_slide(prs.slide_layouts[1])
title = slide_10.shapes.title
content = slide_10.placeholders[1]

title.text = "Attribute Extraction"
content.text = "Analyzed extraction significance.\n\nImportance of early-layer and non-subject representations."

# Slide 11: Conclusion
slide_11 = prs.slides.add_slide(prs.slide_layouts[1])
title = slide_11.shapes.title
content = slide_11.placeholders[1]

title.text = "Conclusion"
content.text = "Summary of key findings.\n\nFuture research directions and limitations."

# Save the presentation
pptx_path = "/mnt/data/Dissecting_Recall_Presentation.pptx"
prs.save(pptx_path)

pptx_path
